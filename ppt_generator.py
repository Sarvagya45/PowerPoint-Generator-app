import os
import re
import json
import requests
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

class PPTGenerator:
    def __init__(self, api_key=None, theme="Classic Light", font_name="Calibri"):
        self.api_key = api_key or os.getenv("GEMINI_API_KEY")
        if not self.api_key:
            raise ValueError("Gemini API Key not found. Please set GEMINI_API_KEY in your .env file.")

        genai.configure(api_key=self.api_key)
        self.model = genai.GenerativeModel("gemini-2.5-flash")
        self.presentation = Presentation()
        
        self.theme = theme
        self.font_name = font_name

    def clean_content(self, text):
        """Convert HTML / markdown returned by Gemini into plain bullet text."""
        if not isinstance(text, str):
            text = str(text)
        
        text = re.sub(r'<br\s*/?>', '\n', text, flags=re.IGNORECASE)
        text = re.sub(r'</li>', '\n', text, flags=re.IGNORECASE)
        text = re.sub(r'<li>', '• ', text, flags=re.IGNORECASE)
        text = re.sub(r'</?ul>', '', text, flags=re.IGNORECASE)
        text = re.sub(r'</?ol>', '', text, flags=re.IGNORECASE)
        text = re.sub(r'</?p>', '\n', text, flags=re.IGNORECASE)
        text = re.sub(r'<[^>]+>', '', text)
        text = text.replace('\\n', '\n')
        text = text.replace('**', '')
        text = re.sub(r'\n{3,}', '\n\n', text)
        return text.strip()

    def generate_content_outline(self, topic, num_slides=5, audience="General"):
        prompt = f"""
        Create a detailed outline for the PowerPoint presentation on "{topic}" with {num_slides} slides.
        The target audience for this presentation is: {audience}.
        Please tailor the tone, complexity, and vocabulary specifically for this audience.
        
        CRITICAL FORMATTING RULE: Limit the content to a maximum of 3 to 4 concise bullet points per slide. 
        Do not write long paragraphs. This ensures the text maintains proper indentation and does not overflow off the slide.
        
        return the response as a JSON array with the following structure:
        [
            {{
                "title": "Slide Title",
                "content": "Main content points as bullet points (max 3-4 points)",
                "slide_type": "title|content|image|conclusion"
            }}
        ]

        Make sure the content is engaging, informative, and well-structured.
        the response must be a valid JSON array.
        """
        try:
            response = self.model.generate_content(prompt)
            content = response.text.strip()

            if "```json" in content:
                content = content.split("```json")[1].split("```")[0].strip()
            elif "```" in content:
                content = content.split("```")[1].strip()

            if not content.startswith("[") or not content.endswith(']'):
                return None

            try:
                return json.loads(content)
            except json.JSONDecodeError as e:
                print(f"JSON Parse Error: {e}")
                return None
        except Exception as e:
            print(f"Generation Error: {e}")
            return None

    def generate_image_description(self, slide_content):
        prompt = f"""
        Based on this slide content, suggest a relevant image description that would enhance the presentation.
        {slide_content}
        Return only a brief, descriptive phrase suitable for image search (max 5 words)
        """
        try:
            response = self.model.generate_content(prompt)
            return response.text.strip()
        except Exception as e:
            return "professional presentation"

    def download_image(self, query, save_path="temp_image.jpg"):
        query_encoded = query.replace(" ", "+")
        try:
            url = "https://api.pexels.com/v1/search"
            pexels_key = os.getenv("PEXELS_API_KEY")
            if pexels_key:
                headers = {"Authorization": pexels_key}
                params = {"query": query, "per_page": 1, "orientation": "landscape"}
                response = requests.get(url, headers=headers, params=params, timeout=10)
                response.raise_for_status()
                data = response.json()
                if data.get('photos'):
                    image_url = data['photos'][0]['src']['large2x']
                    image_response = requests.get(image_url, timeout=15)
                    image_response.raise_for_status()
                    with open(save_path, 'wb') as f:
                        f.write(image_response.content)
                    return save_path
        except Exception as e:
            pass
            
        try:
            image_url = f"https://source.unsplash.com/1280x720/?{query_encoded}"
            image_response = requests.get(image_url, timeout=15, allow_redirects=True)
            image_response.raise_for_status()
            if len(image_response.content) > 5000:  
                with open(save_path, 'wb') as f:
                    f.write(image_response.content)
                return save_path
        except Exception as e:
            pass

        return None

    def _get_theme_colors(self):
        themes = {
            "Classic Light": {"bg": (255, 255, 255), "text": (0, 0, 0), "subtitle": (85, 85, 85)},
            "Midnight Dark": {"bg": (30, 30, 30), "text": (255, 255, 255), "subtitle": (200, 200, 200)},
            "Corporate Blue": {"bg": (15, 76, 117), "text": (255, 255, 255), "subtitle": (200, 220, 240)},
            "Forest Green": {"bg": (27, 67, 50), "text": (255, 255, 255), "subtitle": (180, 210, 190)}
        }
        return themes.get(self.theme, themes["Classic Light"])

    def _apply_theme_to_slide(self, slide):
        colors = self._get_theme_colors()
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*colors["bg"])
        return colors

    def _style_text_frame(self, text_frame, font_size, rgb_tuple, bold=False, align=None):
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = self.font_name
            paragraph.font.size = Pt(font_size)
            paragraph.font.color.rgb = RGBColor(*rgb_tuple)
            if bold:
                paragraph.font.bold = True
            if align:
                paragraph.alignment = align

    def create_title_slide(self, title, subtitle=""):
        slide_layout = self.presentation.slide_layouts[0]
        slide = self.presentation.slides.add_slide(slide_layout)
        colors = self._apply_theme_to_slide(slide)

        title_shape = slide.shapes.title
        title_shape.text = title
        self._style_text_frame(title_shape.text_frame, 32, colors["text"], bold=True, align=PP_ALIGN.CENTER)

        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            self._style_text_frame(subtitle_shape.text_frame, 20, colors["subtitle"], align=PP_ALIGN.CENTER)

    def create_content_slide(self, title, content, include_image=True):
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        colors = self._apply_theme_to_slide(slide)

        title_shape = slide.shapes.title
        title_shape.text = title
        self._style_text_frame(title_shape.text_frame, 30, colors["text"], bold=True)

        content_shape = slide.placeholders[1]
        content_shape.text = content
        
        # Proper spacing bounding boxes to prevent bottom clipping
        if include_image:
            content_shape.left = Inches(0.5)
            content_shape.top = Inches(1.6)
            content_shape.width = Inches(4.5)
            content_shape.height = Inches(5.2)
        else:
            content_shape.left = Inches(0.5)
            content_shape.top = Inches(1.6)
            content_shape.width = Inches(9.0)
            content_shape.height = Inches(5.2)

        self._style_text_frame(content_shape.text_frame, 20, colors["text"])

        if include_image:
            try:
                image_desc = self.generate_image_description(content)
                image_path = self.download_image(image_desc)
                if image_path and os.path.exists(image_path):
                    slide.shapes.add_picture(image_path, Inches(5.2), Inches(1.6), width=Inches(4.3), height=Inches(4.5))
                    os.remove(image_path)
            except Exception as e:
                print(f"Error adding picture: {e}")

        return slide

    def create_image_slide(self, title, image_query, content=""):
        slide_layout = self.presentation.slide_layouts[8] # Title and 2 content often or blank
        slide = self.presentation.slides.add_slide(slide_layout)
        colors = self._apply_theme_to_slide(slide)

        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        title_box.text_frame.text = title
        self._style_text_frame(title_box.text_frame, 28, colors["text"], bold=True, align=PP_ALIGN.CENTER)

        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.5), Inches(5.2))
        content_box.text_frame.text = content
        self._style_text_frame(content_box.text_frame, 20, colors["text"])

        try:
            image_path = self.download_image(image_query)
            if image_path and os.path.exists(image_path):
                slide.shapes.add_picture(image_path, Inches(5.2), Inches(1.6), width=Inches(4.3), height=Inches(4.5))
                os.remove(image_path)
        except Exception as e:
            print(f"Error adding picture: {e}")

    def generate_presentation(self, topic, num_slides=5, output_file="presentation.pptx", progress_callback=None, title_subtitle=""):
        if progress_callback:
            progress_callback(10, "Planning the structure and generating outline with Gemini...")

        content_outline = self.generate_content_outline(topic, num_slides)

        if not content_outline:
            raise Exception("Failed to generate content outline from Gemini. Check API keys and quotas.")

        total_slides = len(content_outline)
        for i, slide_data in enumerate(content_outline):
            title = slide_data.get("title", f"Slide {i+1}")
            content = slide_data.get("content", "")
            if isinstance(content, list):
                content = "\n".join(str(item) for item in content)
            content = self.clean_content(content)
            slide_type = slide_data.get("slide_type", "content")

            if progress_callback:
                progress = 10 + int(((i) / total_slides) * 80)
                progress_callback(progress, f"Generating Slide {i+1} of {total_slides}: {title}")

            if i == 0 or slide_type == "title":
                self.create_title_slide(title, subtitle=title_subtitle if i == 0 else "")
            elif slide_type == "content":
                self.create_content_slide(title, content, include_image=True)
            elif slide_type == "image":
                img_query = self.generate_image_description(content)
                self.create_image_slide(title, img_query, content)
            else:
                include_image = (i % 3 == 0)
                self.create_content_slide(title, content, include_image)

        self.presentation.save(output_file)
        if progress_callback:
             progress_callback(100, f"Presentation completed successfully!")

        return output_file
